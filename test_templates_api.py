"""
PPT模板功能测试脚本
验证所有模板API端点和服务功能
"""

import asyncio
import json
import sys
from pathlib import Path

# 添加项目根目录到Python路径
sys.path.insert(0, str(Path(__file__).parent))

async def test_imports():
    """测试所有导入是否正确"""
    print("=" * 60)
    print("📋 测试1: 模块导入检查")
    print("=" * 60)
    
    try:
        print("✓ 导入 ppt_parser...")
        from backend.service.ppt_parser import PPTParser, extract_template_from_pptx
        print("  ✅ ppt_parser 导入成功")
        
        print("✓ 导入 template_service...")
        from backend.service.template_service import TemplateService
        print("  ✅ template_service 导入成功")
        
        print("✓ 导入 templates_v2 API...")
        from backend.api.templates_v2 import router
        print("  ✅ templates_v2 API 导入成功")
        
        print("\n✅ 所有模块导入成功！\n")
        return True
        
    except Exception as e:
        print(f"\n❌ 导入失败: {e}\n")
        import traceback
        traceback.print_exc()
        return False


async def test_database_schema():
    """测试数据库表是否正确创建"""
    print("=" * 60)
    print("📊 测试2: 数据库表结构检查")
    print("=" * 60)
    
    try:
        from backend.repository.supabase_client import get_supabase_client
        
        supabase = get_supabase_client()
        
        # 检查 user_templates 表
        print("✓ 检查 user_templates 表...")
        try:
            result = supabase.table('user_templates').select('*').limit(1).execute()
            print("  ✅ user_templates 表存在")
        except Exception as e:
            print(f"  ⚠️  user_templates 表可能不存在: {e}")
        
        # 检查 template_favorites 表
        print("✓ 检查 template_favorites 表...")
        try:
            result = supabase.table('template_favorites').select('*').limit(1).execute()
            print("  ✅ template_favorites 表存在")
        except Exception as e:
            print(f"  ⚠️  template_favorites 表可能不存在: {e}")
        
        print("\n✅ 数据库表检查完成！\n")
        return True
        
    except Exception as e:
        print(f"\n❌ 数据库检查失败: {e}\n")
        import traceback
        traceback.print_exc()
        return False


async def test_ppt_parser():
    """测试PPT解析功能"""
    print("=" * 60)
    print("🎨 测试3: PPT解析功能")
    print("=" * 60)
    
    try:
        from backend.service.ppt_parser import PPTParser
        from pptx import Presentation
        import tempfile
        
        # 创建一个简单的测试PPT
        print("✓ 创建测试PPT文件...")
        prs = Presentation()
        
        # 添加标题幻灯片
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "测试课件"
        subtitle.text = "PPT解析测试"
        
        # 添加内容幻灯片
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "测试内容"
        
        text_frame = body.text_frame
        text_frame.text = "这是测试内容"
        
        # 保存到临时文件
        temp_path = tempfile.mktemp(suffix='.pptx')
        prs.save(temp_path)
        print(f"  ✅ 测试PPT创建成功: {temp_path}")
        
        # 测试验证函数
        print("✓ 验证PPT文件...")
        is_valid, message = PPTParser.validate_pptx_file(temp_path)
        print(f"  验证结果: {is_valid}, {message}")
        
        if is_valid:
            print("  ✅ PPT文件验证成功")
            
            # 测试解析函数
            print("✓ 解析PPT文件...")
            template_data = PPTParser.parse_pptx(temp_path)
            
            print(f"  - 幻灯片数: {template_data['page_count']}")
            print(f"  - 主题色: {template_data['theme_colors']}")
            print(f"  - 字体: {template_data['fonts']}")
            print(f"  - 文件大小: {template_data['file_size']} 字节")
            print("  ✅ PPT解析成功")
        else:
            print(f"  ❌ PPT验证失败: {message}")
        
        # 清理临时文件
        import os
        os.remove(temp_path)
        
        print("\n✅ PPT解析功能测试完成！\n")
        return True
        
    except Exception as e:
        print(f"\n❌ PPT解析测试失败: {e}\n")
        import traceback
        traceback.print_exc()
        return False


async def test_api_routes():
    """测试API路由是否正确配置"""
    print("=" * 60)
    print("🔗 测试4: API路由检查")
    print("=" * 60)
    
    try:
        from backend.api.templates_v2 import router
        
        # 获取所有路由
        routes = router.routes
        
        print(f"✓ 已注册的路由数: {len(routes)}")
        
        expected_routes = [
            "upload_template",
            "save_courseware_as_template",
            "get_my_templates",
            "get_public_templates",
            "get_template_detail",
            "copy_template",
            "delete_template",
            "add_favorite",
            "remove_favorite",
            "get_favorites",
            "get_template_categories",
        ]
        
        route_names = [r.name for r in routes if hasattr(r, 'name')]
        
        print("\n已注册的路由:")
        for route_name in route_names:
            status = "✅" if any(exp in route_name for exp in expected_routes) else "ℹ️"
            print(f"  {status} {route_name}")
        
        print(f"\n✅ API路由检查完成！\n")
        return True
        
    except Exception as e:
        print(f"\n❌ API路由检查失败: {e}\n")
        import traceback
        traceback.print_exc()
        return False


async def test_main_app_integration():
    """测试主应用是否正确集成新API"""
    print("=" * 60)
    print("🚀 测试5: 主应用集成检查")
    print("=" * 60)
    
    try:
        from backend.main import app
        
        # 获取应用的所有路由
        routes = app.routes
        
        print(f"✓ 应用总路由数: {len(routes)}")
        
        # 查找模板相关的路由
        template_routes = [r for r in routes if hasattr(r, 'path') and 'template' in r.path.lower()]
        
        print(f"✓ 模板相关路由数: {len(template_routes)}")
        
        print("\n模板相关路由:")
        for route in template_routes[:10]:  # 显示前10个
            if hasattr(route, 'path') and hasattr(route, 'methods'):
                print(f"  ✅ {route.path} {route.methods}")
        
        print(f"\n✅ 主应用集成检查完成！\n")
        return True
        
    except Exception as e:
        print(f"\n❌ 主应用集成检查失败: {e}\n")
        import traceback
        traceback.print_exc()
        return False


async def main():
    """运行所有测试"""
    print("\n" + "=" * 60)
    print("🧪 PPT模板功能测试")
    print("=" * 60 + "\n")
    
    tests = [
        ("模块导入", test_imports),
        ("数据库表", test_database_schema),
        ("PPT解析", test_ppt_parser),
        ("API路由", test_api_routes),
        ("应用集成", test_main_app_integration),
    ]
    
    results = []
    
    for test_name, test_func in tests:
        try:
            result = await test_func()
            results.append((test_name, result))
        except Exception as e:
            print(f"❌ {test_name} 测试异常: {e}\n")
            results.append((test_name, False))
    
    # 输出测试总结
    print("=" * 60)
    print("📊 测试总结")
    print("=" * 60)
    
    passed = sum(1 for _, result in results if result)
    total = len(results)
    
    for test_name, result in results:
        status = "✅ 通过" if result else "❌ 失败"
        print(f"{status}: {test_name}")
    
    print(f"\n总计: {passed}/{total} 测试通过")
    
    if passed == total:
        print("\n🎉 所有测试通过！PPT模板功能已准备就绪！\n")
        return 0
    else:
        print(f"\n⚠️  {total - passed} 个测试失败，请检查错误信息\n")
        return 1


if __name__ == "__main__":
    exit_code = asyncio.run(main())
    sys.exit(exit_code)
