from fastapi import APIRouter, BackgroundTasks, HTTPException
from pydantic import BaseModel
from uuid import uuid4
from pathlib import Path
import json
import os
from core import task_store  # 自定义状态/存储

router = APIRouter(prefix="/api/ppt")

class PPTGenerateRequest(BaseModel):
    user_id: str
    messages: list[dict]   # [{"role":"user","content":"..."}]
    knowledge: list[str]
    slide_count: int = 20
    style: str = "简洁"

class PPTStatus(BaseModel):
    status: str
    message: str = ""
    download_url: str | None = None

def generate_ppt_task(task_id: str, payload: PPTGenerateRequest):
    task_store.set_status(task_id, "running")
    try:
        # 1. 组合 prompt
        prompt = (
            "请根据以下对话与知识库，输出PPT结构：标题/每页大纲，"
            f"页数{payload.slide_count}，风格{payload.style}\n\n"
            "对话:\n" + "\n".join(f"{m['role']}:{m['content']}" for m in payload.messages)
            + "\n\n知识库:\n" + "\n".join(payload.knowledge)
        )
        # 2. 调用大模型（此处假设封装函数）
        from core.llm_client import ask_llm
        outline = ask_llm(prompt, temperature=0.2, max_tokens=1500)
        # 3. 逐页生成内容
        slides = []
        for idx, item in enumerate(parse_outline(outline, payload.slide_count), 1):
            text_prompt = (
                f"请为PPT第{idx}页生成内容，页标题：{item['title']}\n"
                + "关键点："+";".join(item["points"])
            )
            slide_text = ask_llm(text_prompt, temperature=0.3, max_tokens=800)
            slides.append({"title": item["title"], "content": slide_text})

        # 4. 生成 pptx
        from pptx import Presentation
        prs = Presentation()
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.text = s["title"]
            body = slide.shapes.placeholders[1]
            body.text = s["content"]

        out_dir = Path("storage/ppt")
        out_dir.mkdir(parents=True, exist_ok=True)
        out_file = out_dir / f"{task_id}.pptx"
        prs.save(str(out_file))

        task_store.set_status(task_id, "finished", download_url=f"/api/ppt/download/{task_id}")
    except Exception as e:
        task_store.set_status(task_id, "failed", message=str(e))

def parse_outline(outline_text: str, count: int):
    # 简单占位解析，真实需稳健
    lines = [l.strip() for l in outline_text.splitlines() if l.strip()]
    items = []
    for l in lines[:count]:
        items.append({"title": l, "points": ["自动提炼要点"]})
    if not items:
        items.append({"title": "自动生成内容", "points": ["..."]})
    return items

@router.post("/generate")
def create_ppt(payload: PPTGenerateRequest, background_tasks: BackgroundTasks):
    task_id = str(uuid4())
    task_store.set_status(task_id, "pending")
    background_tasks.add_task(generate_ppt_task, task_id, payload)
    return {"task_id": task_id}

@router.get("/status/{task_id}", response_model=PPTStatus)
def get_status(task_id: str):
    info = task_store.get_status(task_id)
    if not info:
        raise HTTPException(status_code=404, detail="task not found")
    return info

@router.get("/download/{task_id}")
def download_ppt(task_id: str):
    info = task_store.get_status(task_id)
    if not info or info["status"] != "finished":
        raise HTTPException(status_code=404, detail="no file")
    filepath = Path("storage/ppt") / f"{task_id}.pptx"
    return FileResponse(str(filepath), filename=f"{task_id}.pptx")