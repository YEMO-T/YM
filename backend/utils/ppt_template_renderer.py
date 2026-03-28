import os
import io
import logging
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from typing import List, Dict, Any, Optional
from backend.schema.chat_schema import PPTSlide

logger = logging.getLogger(__name__)

class PPTTemplateRenderer:
    """
    基于模板的 PPT 渲染引擎。
    负责将 AI 生成的内容注入到现有的 .pptx 模板中。
    """

    def __init__(self, template_path: str):
        if not os.path.exists(template_path):
            logger.error(f"模板文件不存在: {template_path}")
            raise FileNotFoundError(f"Template not found at {template_path}")
        
        self.prs = Presentation(template_path)
        self.layouts = self.prs.slide_layouts
        logger.info(f"成功加载模板: {template_path}, 包含 {len(self.layouts)} 个布局")

    def _find_layout_by_name(self, keywords: List[str]) -> Any:
        """根据名称关键词查找最匹配的布局"""
        for layout in self.layouts:
            name_lower = layout.name.lower()
            if any(k.lower() in name_lower for k in keywords):
                return layout
        return None

    def get_layout_for_type(self, page_type: str) -> Any:
        """根据页面类型选择布局"""
        page_type = page_type.lower()
        
        if page_type == "cover":
            # 优先找名字含 cover/title 的，找不到用第0个
            return self._find_layout_by_name(["cover", "title", "封面"]) or self.layouts[0]
        elif page_type == "toc":
            return self._find_layout_by_name(["toc", "agenda", "content", "目录"]) or self.layouts[1]
        elif page_type == "ending":
            return self._find_layout_by_name(["ending", "thanks", "thank", "结束", "致谢"]) or self.layouts[0]
        else:
            # 普通内容页
            return self._find_layout_by_name(["content", "body", "normal", "内容", "正文"]) or self.layouts[1]

    def fill_placeholders(self, slide, title: str, content: List[str]):
        """填充占位符并保留样式"""
        # 填充标题
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # 寻找正文占位符 (通常是 Body 或序号较大的 Placeholder)
        body_placeholder = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx > 0 and shape.has_text_frame:
                body_placeholder = shape
                break
        
        if body_placeholder:
            tf = body_placeholder.text_frame
            # 保持原有样式，只清空文本
            # 如果是列表，我们逐项添加
            tf.clear() # 这会清空所有段落，但保留文本框样式
            
            for i, text in enumerate(content):
                p = tf.add_paragraph()
                p.text = str(text)
                p.level = 0
                # 如果是目录页，可能需要特殊对齐
                # p.alignment = PP_ALIGN.LEFT

    def render(self, slides_data: List[PPTSlide]) -> io.BytesIO:
        """全量渲染 PPT 并返回内存流"""
        # 先删除模板中原有的所有幻灯片（如果是一个空母版，可能本来就没有）
        # 注意：python-pptx 不直接支持删除幻灯片，但我们可以从空白开始
        # 或者我们假设模板就是一个包含几个示例页的 PPT，我们需要清理后再加
        
        # 技巧：如果是纯模板，通常 prs.slides 是空的
        # 如果不是空的，我们可能需要保留母版而删除实例页
        
        # 简单的做法：直接 add_slide，然后如果原本有页，最后手工处理（通常模板文件就是空的页）
        
        for slide_data in slides_data:
            layout = self.get_layout_for_type(slide_data.page_type)
            slide = self.prs.slides.add_slide(layout)
            self.fill_placeholders(slide, slide_data.title, slide_data.content)
            
        # 导出为二进制流
        pptx_io = io.BytesIO()
        self.prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

def render_ppt_with_template(slides: List[PPTSlide], template_id: str) -> io.BytesIO:
    """高层封装函数"""
    # 路径匹配 TemplateService.upload_template 中的 target_path
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    template_path = os.path.join(base_dir, 'data', 'templates', f"{template_id}.pptx")
    
    # 容错：如果找不到特定模板，尝试用 default.pptx 或报错
    if not os.path.exists(template_path):
        logger.warning(f"模板 {template_id} 不存在于本地，尝试查找默认模板")
        template_path = os.path.join(base_dir, 'data', 'templates', "default.pptx")
        if not os.path.exists(template_path):
             # 最后保底：直接用 python-pptx 默认生成（不带样式的）
             from .ppt_generator import generate_pptx
             return generate_pptx(slides)

    renderer = PPTTemplateRenderer(template_path)
    return renderer.render(slides)
