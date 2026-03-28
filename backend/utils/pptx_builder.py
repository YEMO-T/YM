import io
import os
import logging
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Inches, Pt
from typing import List, Dict, Optional

logger = logging.getLogger(__name__)

class PPTXBuilder:
    """
    基于 python-pptx 的专业 PPTX 构建器。
    用于将 AI 内容填入物理模板中，保留母版设计。
    """

    def __init__(self, template_path: Optional[str] = None):
        self.template_path = template_path
        try:
            if template_path and os.path.exists(template_path):
                self.prs = Presentation(template_path)
                # 清除模板中的演示页面，仅保留母版
                self._clear_slides()
                logger.info(f"✓ PPTXBuilder 已载入模板: {os.path.basename(template_path)}")
            else:
                self.prs = Presentation()
                logger.info("✓ PPTXBuilder 已载入默认空白。")
        except Exception as e:
            logger.error(f"⚠ PPTXBuilder 初始化失败: {e}")
            self.prs = Presentation()

    def _clear_slides(self):
        """移除物理模板中已有的幻灯片内容。"""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def _find_suitable_layout(self, slide_type: str = 'content'):
        """
        根据幻灯片类型寻找最接近的布局。
        """
        layouts = self.prs.slide_layouts
        
        # 常见布局名称映射 (中英文支持)
        layout_map = {
            'cover': ['Title Slide', '封面', '标题', 'Title'],
            'content': ['Title and Content', '正文', '标题和内容', 'Content'],
            'summary': ['Section Header', '总结', '结束', 'Summary', 'Closing'],
        }
        
        target_names = layout_map.get(slide_type, layout_map['content'])
        
        # 1. 尝试按名称匹配
        for i, layout in enumerate(layouts):
            for name in target_names:
                if name.lower() in layout.name.lower():
                    return i
        
        # 2. 回退机制：封面通常是 index 0，正文通常是 index 1
        if slide_type == 'cover': return 0
        if len(layouts) > 1: return 1
        return 0

    def add_slide_from_data(self, slide_data: Dict):
        """
        根据单页 JSON 数据添加幻灯片。
        """
        stype = slide_data.get('type', 'content').lower()
        layout_idx = self._find_suitable_layout(stype)
        slide_layout = self.prs.slide_layouts[layout_idx]
        
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 填充标题 (Shape.title 通常是 ID=1)
        title_text = slide_data.get('title', '')
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        
        # 填充内容 (遍历占位符)
        content_text = slide_data.get('content', '')
        if content_text:
            # 找到正文占位符 (优先使用 BODY 类型的占位符)
            body_ph = None
            for shape in slide.placeholders:
                if shape == slide.shapes.title:
                    continue
                # 检查是否为正文占位符
                ph_fmt = getattr(shape, 'placeholder_format', None)
                if ph_fmt and ph_fmt.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    body_ph = shape
                    break
            
            # 如果没找到特定的 BODY，则尝试任何有文本框的占位符
            if not body_ph:
                for shape in slide.placeholders:
                    if shape != slide.shapes.title and hasattr(shape, 'text_frame'):
                        body_ph = shape
                        break
            
            if body_ph:
                body_ph.text = content_text

        # 尝试处理图片占位符 (如果有 imagePrompt)
        # NOTE: 暂不真正下载图片，仅作为逻辑预留
        image_prompt = slide_data.get('imagePrompt')
        if image_prompt:
            logger.debug(f"Slide '{title_text}' requires image: {image_prompt}")
            #TODO: 调用图片生成 API 并插入到图片占位符或背景
        
    def build_stream(self) -> io.BytesIO:
        """生成内存流。"""
        stream = io.BytesIO()
        self.prs.save(stream)
        stream.seek(0)
        return stream

def create_pptx_from_ai_json(slides_data: List[Dict], template_id: Optional[str] = None) -> io.BytesIO:
    """
    高层封装：直接将 JSON 的 slides 数组转为 PPTX 流。
    """
    template_path = None
    if template_id:
        # 指向 data/templates 目录
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        template_path = os.path.join(base_dir, 'data', 'templates', f"{template_id}.pptx")

    builder = PPTXBuilder(template_path)
    for slide_data in slides_data:
        builder.add_slide_from_data(slide_data)
        
    return builder.build_stream()
