import os
import io
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Any, Optional
from backend.schema.chat_schema import PPTSlide

logger = logging.getLogger(__name__)

class PPTGenerator:
    """
    市场级 PPT 生成引擎。
    目标：零报错、自动排版、精美视觉。
    """
    
    def __init__(self, template_path: Optional[str] = None):
        try:
            if template_path and os.path.exists(template_path):
                self.prs = Presentation(template_path)
                logger.info(f"使用模板生成 PPT: {template_path}")
            else:
                self.prs = Presentation()
                logger.info("使用默认母版生成 PPT")
        except Exception as e:
            logger.error(f"初始化 Presentation 失败: {e}")
            self.prs = Presentation() # 强制回退到默认

    def _set_font(self, run, size_pt: int = 18, color_rgb: tuple = (0, 0, 0), bold: bool = False):
        """统一字体设置"""
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color_rgb)
        # 尝试设置中文通用字体
        try:
            run.font.name = '微软雅黑'
            run.element.rPr.set('asciiFont', '微软雅黑')
            run.element.rPr.set('hAnsiFont', '微软雅黑')
            run.element.rPr.set('eastAsiaFont', '微软雅黑')
        except:
            run.font.name = 'Arial'

    def add_cover(self, title: str, subtitle: str = ""):
        """渲染封面页"""
        layout = self.prs.slide_layouts[0] # Title Slide
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

    def add_toc(self, title: str, items: List[str]):
        """渲染目录页"""
        layout = self.prs.slide_layouts[1] # Title and Content
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear() # 清空默认
            for i, item in enumerate(items):
                p = tf.add_paragraph()
                p.text = f"{i+1}. {item}"
                p.level = 0
                p.space_after = Pt(12)

    def add_content(self, title: str, bullets: List[str]):
        """渲染内容页"""
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True
            for point in bullets:
                p = tf.add_paragraph()
                p.text = str(point)
                p.level = 0
                p.space_before = Pt(6)

    def add_ending(self, title: str = "感谢观看", subtitle: str = "Q&A"):
        """渲染结束页"""
        # 有些模板没有专门结束页，我们用 Title Slide 或简单 Content 模拟
        layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def generate(self, slides_data: List[PPTSlide]) -> io.BytesIO:
        """执行全量渲染"""
        for i, slide in enumerate(slides_data):
            try:
                stype = slide.page_type.lower()
                if stype == "cover" or i == 0:
                    subtitle = slide.content[0] if slide.content else "PPT 自动生成"
                    self.add_cover(slide.title, subtitle)
                elif stype == "toc":
                    self.add_toc(slide.title, slide.content)
                elif stype == "ending" or i == len(slides_data) - 1:
                    self.add_ending(slide.title)
                else:
                    self.add_content(slide.title, slide.content)
            except Exception as e:
                logger.error(f"渲染第 {i} 页失败: {e}")
                # 继续尝试渲染下一页，不要因为一页崩溃全量
                continue
                
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output

def generate_pptx(slides: List[PPTSlide]) -> io.BytesIO:
    """快速渲染接口"""
    engine = PPTGenerator()
    return engine.generate(slides)
