"""
PPT文件解析器 - 提取PPT结构、样式、占位符等信息
支持 .pptx 格式（Microsoft Office Open XML）
"""

import io
import base64
import logging
from typing import Dict, List, Any, Optional, Tuple
from pathlib import Path
from PIL import Image
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PPTParser:
    """PPT文件解析器 - 解析PPTX文件结构和样式信息"""
    
    # 常见占位符类型映射
    PLACEHOLDER_TYPES = {
        'title': 'Title',
        'body': 'Body',
        'subtitle': 'Subtitle',
        'table': 'Table',
        'chart': 'Chart',
        'image': 'Image',
        'footer': 'Footer',
    }
    
    @staticmethod
    def parse_pptx(file_path: str) -> Dict[str, Any]:
        """
        解析PPT文件并提取关键信息
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            {
                'slides_structure': [...],      # 幻灯片版式信息
                'theme_colors': {...},          # 配色方案
                'fonts': {...},                 # 字体配置
                'placeholders': {...},          # 占位符定义
                'thumbnail': base64_string,     # 缩略图
                'page_count': int,              # 页数
                'file_size': int,               # 文件大小（字节）
            }
        """
        try:
            prs = Presentation(file_path)
            
            logger.info(f"开始解析PPT文件: {file_path}")
            
            # 提取各类信息
            slides_structure = PPTParser._extract_slides_structure(prs)
            theme_colors = PPTParser._extract_theme_colors(prs)
            fonts = PPTParser._extract_fonts(prs)
            placeholders = PPTParser._extract_placeholders(prs)
            thumbnail = PPTParser._generate_thumbnail(prs)
            
            # 获取文件大小
            file_size = Path(file_path).stat().st_size if Path(file_path).exists() else 0
            
            result = {
                'slides_structure': slides_structure,
                'theme_colors': theme_colors,
                'fonts': fonts,
                'placeholders': placeholders,
                'thumbnail': thumbnail,
                'page_count': len(prs.slides),
                'file_size': file_size,
                'slide_dimensions': {
                    'width': float(prs.slide_width),
                    'height': float(prs.slide_height),
                },
            }
            
            logger.info(f"✓ PPT解析完成: {len(prs.slides)} 页，{theme_colors.get('primary', 'N/A')} 主色调")
            return result
            
        except Exception as e:
            logger.error(f"PPT解析失败: {type(e).__name__}: {e}")
            raise
    
    @staticmethod
    def _extract_slides_structure(prs: Presentation) -> List[Dict[str, Any]]:
        """提取幻灯片版式信息"""
        slides_info = []
        
        for idx, slide in enumerate(prs.slides):
            slide_data = {
                'index': idx,
                'layout_name': slide.slide_layout.name,
                'shapes': [],
                'background_color': None,
            }
            
            # 提取背景色与背景图
            try:
                if hasattr(slide, 'background') and slide.background:
                    fill = slide.background.fill
                    if fill.type:
                        bg_color = PPTParser._extract_color(fill)
                        slide_data['background_color'] = bg_color
                    
                    # 尝试提取背景图片 (如果存在)
                    # NOTE: python-pptx 对背景图片的直接访问有限，通常作为 shape 处理或在 slide_master 中
            except Exception as e:
                logger.debug(f"背景提取失败: {e}")
            
            # 提取形状和文本框
            for shape in slide.shapes:
                shape_info = {
                    'shape_type': str(shape.shape_type),
                    'name': shape.name,
                    'left': float(shape.left) if hasattr(shape, 'left') else None,
                    'top': float(shape.top) if hasattr(shape, 'top') else None,
                    'width': float(shape.width) if hasattr(shape, 'width') else None,
                    'height': float(shape.height) if hasattr(shape, 'height') else None,
                }
                
                # 提取形状填充色
                if hasattr(shape, 'fill') and shape.fill.type:
                    try:
                        shape_info['fill_color'] = PPTParser._extract_color(shape.fill)
                    except:
                        pass

                # 提取文本内容
                if hasattr(shape, 'text'):
                    shape_info['text'] = shape.text[:200]
                    shape_info['text_length'] = len(shape.text)
                
                # 提取图片内容 (如果是图片)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        base64_img = base64.b64encode(image_bytes).decode('utf-8')
                        shape_info['image_data'] = f"data:{image.ext};base64,{base64_img}"
                    except Exception as e:
                        logger.debug(f"图片提取失败: {e}")

                # 提取文本格式
                if hasattr(shape, 'text_frame'):
                    shape_info['text_format'] = PPTParser._extract_text_format(shape.text_frame)
                
                slide_data['shapes'].append(shape_info)
            
            slides_info.append(slide_data)
        
        logger.info(f"✓ 提取 {len(slides_info)} 个幻灯片结构")
        return slides_info
    
    @staticmethod
    def _extract_theme_colors(prs: Presentation) -> Dict[str, str]:
        """提取PPT主题色"""
        colors = {
            'primary': '#000000',
            'secondary': '#FFFFFF',
            'accent1': '#CCCCCC',
            'accent2': '#CCCCCC',
            'accent3': '#CCCCCC',
        }
        
        try:
            for i, slide in enumerate(prs.slides):
                if i >= 5: break
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'fill') and shape.fill.type:
                            rgb = PPTParser._extract_color(shape.fill)
                            if rgb and rgb not in ['#FFFFFF', '#000000']:
                                colors['primary'] = rgb
                                break
                    except Exception:
                        pass
                    
                    try:
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if hasattr(run.font, 'color') and run.font.color.type:
                                        rgb = PPTParser._extract_color(run.font.color)
                                        if rgb:
                                            colors['primary'] = rgb
                                            break
                    except Exception:
                        pass
        except Exception as e:
            logger.warning(f"主题色提取失败: {e}")
        
        logger.info(f"✓ 提取主题色: 主色调={colors['primary']}")
        return colors
    
    @staticmethod
    def _extract_fonts(prs: Presentation) -> Dict[str, Any]:
        """提取字体信息"""
        fonts_found = {
            'titles': set(),
            'body': set(),
            'sizes': set(),
        }
        
        try:
            for i, slide in enumerate(prs.slides):
                if i >= 10: break
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if hasattr(run.font, 'name') and run.font.name:
                                        fonts_found['body'].add(run.font.name)
                                    if hasattr(run.font, 'size') and run.font.size:
                                        fonts_found['sizes'].add(float(run.font.size))
                    except Exception:
                        pass
        except Exception as e:
            logger.warning(f"字体提取失败: {e}")
        
        # 转换为列表并规避 lint 错误
        body_fonts = list(fonts_found['body'])
        font_sizes = sorted(list(fonts_found['sizes']))
        
        result = {
            'common_fonts': body_fonts[:3] if body_fonts else ['Arial', 'Calibri', 'SimHei'],
            'font_sizes': font_sizes[:5] if font_sizes else [18, 24, 32, 44],
            'default_font': body_fonts[0] if body_fonts else 'Arial',
        }
        
        logger.info(f"✓ 提取字体: {result['default_font']}")
        return result
    
    @staticmethod
    def _extract_placeholders(prs: Presentation) -> Dict[str, Any]:
        """提取占位符定义"""
        placeholders_found = {
            'title': False,
            'content': False,
            'subtitle': False,
            'footer': False,
            'example_content': [],
        }
        
        try:
            for i, slide in enumerate(prs.slides):
                if i >= 5: break
                for shape in slide.shapes:
                    try:
                        # 检查是否为占位符
                        if shape.is_placeholder:
                            ph = shape.placeholder_format
                            ph_type = ph.type
                            
                            if 'title' in str(ph_type).lower():
                                placeholders_found['title'] = True
                            elif 'body' in str(ph_type).lower():
                                placeholders_found['content'] = True
                            elif 'subtitle' in str(ph_type).lower():
                                placeholders_found['subtitle'] = True
                        
                        # 收集示例内容
                        if hasattr(shape, 'text') and shape.text and len(shape.text) > 3:
                            example_content = placeholders_found.get('example_content', [])
                            if isinstance(example_content, list):
                                example_content.append({
                                    'type': shape.name,
                                    'text': shape.text[:50],
                                })
                                placeholders_found['example_content'] = example_content
                    except Exception:
                        pass
        except Exception as e:
            logger.warning(f"占位符提取失败: {e}")
        
        # 限制示例数量
        example_content = placeholders_found.get('example_content', [])
        placeholders_found['example_content'] = example_content[:3] if example_content else []
        
        logger.info(f"✓ 提取占位符: 标题={placeholders_found['title']}, "
                   f"内容={placeholders_found['content']}")
        return placeholders_found
    
    @staticmethod
    def _extract_text_format(text_frame) -> Dict[str, Any]:
        """提取文本格式信息"""
        format_info = {
            'alignment': None,
            'font_size': None,
            'font_name': None,
            'bold': False,
            'italic': False,
            'color': None,
        }
        
        try:
            if text_frame.paragraphs:
                para = text_frame.paragraphs[0]
                
                # 段落对齐
                if hasattr(para, 'alignment') and para.alignment:
                    format_info['alignment'] = str(para.alignment)
                
                # 字体信息
                if para.runs:
                    run = para.runs[0]
                    if hasattr(run.font, 'name'):
                        format_info['font_name'] = run.font.name
                    if hasattr(run.font, 'size') and run.font.size:
                        format_info['font_size'] = float(run.font.size)
                    if hasattr(run.font, 'bold'):
                        format_info['bold'] = run.font.bold
                    if hasattr(run.font, 'italic'):
                        format_info['italic'] = run.font.italic
                    if hasattr(run.font, 'color'):
                        try:
                            format_info['color'] = PPTParser._extract_color(run.font.color)
                        except:
                            pass
        except Exception as e:
            logger.debug(f"文本格式提取部分失败: {e}")
        
        return format_info
    
    @staticmethod
    def _extract_color(color_obj) -> Optional[str]:
        """
        提取颜色信息
        返回十六进制颜色字符串，如 '#FF0000'
        """
        try:
            if hasattr(color_obj, 'rgb'):
                rgb = color_obj.rgb
                if rgb:
                    r = rgb[0]
                    g = rgb[1]
                    b = rgb[2]
                    return f'#{r:02X}{g:02X}{b:02X}'
            
            if hasattr(color_obj, 'type') and str(color_obj.type) == 'RGB':
                rgb_str = str(color_obj).split('RGB')[1]
                return rgb_str
        except:
            pass
        
        return None
    
    @staticmethod
    def _generate_thumbnail(prs: Presentation, max_size: Tuple[int, int] = (300, 225)) -> str:
        """
        生成PPT缩略图（第一页）
        返回Base64编码的PNG图像
        
        Args:
            prs: Presentation对象
            max_size: 缩略图最大尺寸
            
        Returns:
            Base64编码的图像字符串
        """
        try:
            if not prs.slides:
                logger.warning("PPT无幻灯片，无法生成缩略图")
                return ""
            
            # 使用LibreOffice或其他工具转换（这里使用简化方案）
            # 实际环境中可能需要使用 libreoffice --headless 命令
            # 或者使用其他库如 python-pptx-image 等
            
            slide = prs.slides[0]
            
            # 简化方案：使用PIL创建一个带有幻灯片信息的图像
            # 实际项目中应该使用专业的PPT转图片库
            img = Image.new('RGB', max_size, color='white')
            
            # 在这里可以添加更复杂的缩略图生成逻辑
            # 当前为简化版本
            
            # 转换为Base64
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            buffer.seek(0)
            base64_str = base64.b64encode(buffer.getvalue()).decode('utf-8')
            
            logger.info(f"✓ 生成缩略图: {len(base64_str)} 字节")
            return f"data:image/png;base64,{base64_str}"
            
        except Exception as e:
            logger.warning(f"缩略图生成失败: {e}")
            return ""
    
    @staticmethod
    def validate_pptx_file(file_path: str) -> Tuple[bool, str]:
        """
        验证PPT文件有效性
        
        Returns:
            (is_valid, error_message)
        """
        try:
            if not Path(file_path).exists():
                return False, "文件不存在"
            
            if not file_path.lower().endswith('.pptx'):
                return False, "只支持 .pptx 格式"
            
            # 尝试打开文件
            prs = Presentation(file_path)
            
            if not prs.slides:
                return False, "PPT文件为空（无幻灯片）"
            
            return True, "验证成功"
            
        except Exception as e:
            return False, f"文件验证失败: {str(e)}"


def extract_template_from_pptx(file_path: str) -> Dict[str, Any]:
    """
    从PPT文件提取模板信息
    
    Args:
        file_path: PPT文件路径
        
    Returns:
        模板信息字典，包含所有结构和样式信息
    """
    # 验证文件
    is_valid, message = PPTParser.validate_pptx_file(file_path)
    if not is_valid:
        raise ValueError(f"PPT验证失败: {message}")
    
    # 解析PPT
    template_data = PPTParser.parse_pptx(file_path)
    
    return template_data
